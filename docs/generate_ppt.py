import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def parse_markdown(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
        
    slides = []
    current_slide = None
    state = "NONE" # NONE, CORE_TEXT, SCRIPT
    
    # 슬라이드 감지를 위한 정규식 (#### 슬라이드 N: [제목] 혹은 #### 슬라이드 N [제목])
    slide_pattern = re.compile(r'^#### 슬라이드 (\d+):?\s*(.*)')
    
    for line in lines:
        stripped = line.strip()
        
        # 새 슬라이드 시작 감지
        match = slide_pattern.match(stripped)
        if match:
            if current_slide:
                slides.append(current_slide)
            slide_num = int(match.group(1))
            current_slide = {
                'num': slide_num,
                'title': match.group(2).strip(),
                'design': "",
                'bullets': [],
                'script_lines': []
            }
            state = "NONE"
            continue
            
        if not current_slide:
            continue
            
        # 속성별 파싱 분기
        if stripped.startswith('*   **슬라이드 제목**:'):
            current_slide['title'] = stripped.replace('*   **슬라이드 제목**:', '').strip()
            state = "NONE"
        elif stripped.startswith('*   **시각 디자인**:'):
            current_slide['design'] = stripped.replace('*   **시각 디자인**:', '').strip()
            state = "NONE"
        elif stripped.startswith('*   **핵심 텍스트**:'):
            state = "CORE_TEXT"
        elif stripped.startswith('*   **발표 대본 (Speaker Note)**:'):
            state = "SCRIPT"
        else:
            if state == "CORE_TEXT":
                # 핵심 텍스트 라인 파싱 (불릿 처리)
                m = re.match(r'^\s*[\*\-\+➔]\s*(.*)', line)
                if m:
                    content = m.group(1).strip()
                    # 들여쓰기 감지 (4칸 이상 들여쓰기 시 하위 수준 불릿)
                    indent_match = re.match(r'^(\s*)', line)
                    indent_len = len(indent_match.group(1)) if indent_match else 0
                    if indent_len >= 4:
                        current_slide['bullets'].append("  - " + content)
                    else:
                        current_slide['bullets'].append(content)
                elif stripped:
                    current_slide['bullets'].append(stripped)
            elif state == "SCRIPT":
                # 발표 대본 라인 파싱
                if stripped.startswith('>'):
                    script_content = stripped[1:].strip()
                    # 앞뒤 따옴표 제거
                    if script_content.startswith('"') or script_content.startswith('“'):
                        script_content = script_content[1:]
                    if script_content.endswith('"') or script_content.endswith('”'):
                        script_content = script_content[:-1]
                    if script_content:
                        current_slide['script_lines'].append(script_content)
                        
    if current_slide:
        slides.append(current_slide)
        
    # 발표 대본 정제
    for slide in slides:
        slide['script'] = "\n".join(slide['script_lines']).replace('<br/>', '\n').replace('\\n', '\n')
        del slide['script_lines']
        
    return slides

def create_ppt(slides_data, output_path):
    prs = Presentation()
    # 16:9 와이드스크린 해상도 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 현대적인 어두운 테마(Dark Premium Style) 색상 토큰 정의
    DARK_BG = RGBColor(26, 26, 30)      # 짙은 차콜/블랙 배경
    WHITE = RGBColor(245, 245, 247)     # 크림 화이트 (텍스트)
    ACCENT = RGBColor(255, 107, 107)    # 코랄 레드 (강조 및 제목)
    LIGHT_GRAY = RGBColor(180, 180, 185) # 소프트 그레이 (서브 불릿)

    # 레이아웃 정의 (6번: 빈 화면 레이아웃)
    blank_layout = prs.slide_layouts[6]
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # 1. 배경 설정 (단색 채우기)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        # 2. 제목 텍스트 박스 추가 (좌표: L=1.0", T=0.8", W=11.333", H=1.0")
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        
        p_title = tf_title.paragraphs[0]
        p_title.text = f"Slide {slide_info['num']}: {slide_info['title']}"
        p_title.font.name = '맑은 고딕'
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = ACCENT
        
        # 3. 본문 텍스트 박스 추가 (좌표: L=1.0", T=2.2", W=11.333", H=4.5")
        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.5))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        tf_body.margin_left = tf_body.margin_top = tf_body.margin_right = tf_body.margin_bottom = 0
        
        # 핵심 텍스트 리스트가 비어 있으면 visual guide를 대신 넣어 줌
        bullets = slide_info['bullets'] if slide_info['bullets'] else [f"[Visual Layout Guide] {slide_info['design']}"]
        
        for idx, bullet in enumerate(bullets):
            if idx == 0:
                p = tf_body.paragraphs[0]
            else:
                p = tf_body.add_paragraph()
            
            p.font.name = '맑은 고딕'
            
            # 서브 불릿 레벨 설정
            if bullet.startswith("  -"):
                p.text = bullet.replace("  -", "").strip()
                p.level = 1
                p.font.size = Pt(15)
                p.font.color.rgb = LIGHT_GRAY
            else:
                p.text = bullet.strip()
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = WHITE
                
            p.space_after = Pt(10)

        # 4. 발표자 대본(Speaker Notes) 등록
        if slide_info['script']:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_info['script']
            
    prs.save(output_path)
    print(f"PowerPoint generation completed. File saved at: {output_path}")

if __name__ == "__main__":
    md_file_path = "./docs/9_ppt_presentation_slides.md"
    ppt_output_path = "./docs/Cobot_Fire_Sim_Presentation.pptx"
    
    if os.path.exists(md_file_path):
        print("Parsing markdown slide content...")
        slides = parse_markdown(md_file_path)
        print(f"Parsed {len(slides)} slides.")
        
        print("Building PowerPoint (.pptx) file...")
        create_ppt(slides, ppt_output_path)
    else:
        print(f"Error: {md_file_path} not found.")
